import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Absolute Image Paths (Generated Assets)
IMAGE_COVER = r"C:\Users\karthiksuresh\.gemini\antigravity-ide\brain\a98728d7-c389-4abd-8788-9d3351ddff4c\presentation_cover_1780248577329.png"
IMAGE_DASHBOARD = r"C:\Users\karthiksuresh\.gemini\antigravity-ide\brain\a98728d7-c389-4abd-8788-9d3351ddff4c\dashboard_mockup_1780248595648.png"
IMAGE_WORKFLOW = r"C:\Users\karthiksuresh\.gemini\antigravity-ide\brain\a98728d7-c389-4abd-8788-9d3351ddff4c\workflow_automation_1780248614030.png"

# Color Palette (Corporate HSL Tailored)
COLOR_NAVY = RGBColor(12, 20, 36)
COLOR_GOLD = RGBColor(181, 139, 42)
COLOR_PRIMARY = RGBColor(173, 46, 36) # Red accent
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_LIGHT_GRAY = RGBColor(245, 247, 251)
COLOR_DARK_GRAY = RGBColor(16, 24, 40)
COLOR_MUTED = RGBColor(102, 112, 133)

def add_title(slide, text, color=COLOR_DARK_GRAY, size=32, top=0.6, left=0.8):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(11), Inches(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.bold = True
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.name = "Arial"
    return top + 0.9

def add_bullets(slide, items, top=1.8, left=0.8, width=6.5, height=5, size=15, color=COLOR_DARK_GRAY):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        # Check if item is an empty line for spacing
        if item == "":
            p.text = ""
            continue
            
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Arial"
        
        # Style headings within bullets
        if ":" in item and len(item.split(":")[0]) < 25:
            # Simple bolding of heading prefix
            p.level = 0
            p.font.bold = True
            p.font.size = Pt(size + 1)
        else:
            p.level = 0
            p.font.bold = False
            
    return top + height

def main():
    print("Initializing PPTX Generation...")
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # 16:9 widescreen
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    
    # -------------------------------------------------------------
    # SLIDE 1: COVER SLIDE
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    if os.path.exists(IMAGE_COVER):
        slide1.shapes.add_picture(IMAGE_COVER, 0, 0, width=Inches(13.33), height=Inches(7.5))
    
    # Text Overlay on Cover
    txBox = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p_badge = tf.paragraphs[0]
    p_badge.text = "BUSINESS SYSTEMS ENGINEERING"
    p_badge.font.size = Pt(14)
    p_badge.font.bold = True
    p_badge.font.color.rgb = COLOR_GOLD
    p_badge.alignment = PP_ALIGN.LEFT
    
    p_title = tf.add_paragraph()
    p_title.text = "EAGLE INNOVATIONS"
    p_title.font.size = Pt(56)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_WHITE
    p_title.alignment = PP_ALIGN.LEFT
    
    p_subtitle = tf.add_paragraph()
    p_subtitle.text = "High-Performance CRMs, Active Dashboards, and Process Automations for Growing Teams."
    p_subtitle.font.size = Pt(18)
    p_subtitle.font.color.rgb = COLOR_LIGHT_GRAY
    p_subtitle.alignment = PP_ALIGN.LEFT
    
    # -------------------------------------------------------------
    # SLIDE 2: THE OPERATIONAL CHALLENGE
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    add_title(slide2, "The Core Challenge: Software That Doesn't Fit")
    
    bullets_s2 = [
        "Rigid Off-The-Shelf Solutions: Buying standardized portals or CRMs forces your operational teams to change how they work, creating inefficiencies.",
        "",
        "Spreadsheet & Manual Chaos: Growing teams often default to copy-pasting customer data across scattered spreadsheets, causing human errors and data silos.",
        "",
        "Fragmented Vision: Without a single, central operational layer, leadership loses visibility into real-time metrics, pipeline bottlenecks, and team performance.",
        "",
        "Our Solution: Eagle Innovations engineers custom business software designed specifically around your exact processes, database models, and executive workflows."
    ]
    add_bullets(slide2, bullets_s2, width=11.5)
    
    # -------------------------------------------------------------
    # SLIDE 3: CAPABILITIES (SPLIT)
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    add_title(slide3, "Core Capabilities: Digital Systems Studio")
    
    bullets_s3 = [
        "1. Custom CRMs & Client Portals",
        "We build centralized hubs to manage customer communication, leads, ticket flows, and operations in a unified workspace.",
        "",
        "2. Executive BI Dashboards",
        "Consolidate multiple data streams into real-time analytical command centers, allowing leadership to spot performance trends and operational risks.",
        "",
        "3. Workflow & AI Automation",
        "Connect forms, approvals, cloud storage, notifications, and reports to save your staff 10+ hours per week per employee by automating repetitive work."
    ]
    add_bullets(slide3, bullets_s3, width=6.2)
    
    if os.path.exists(IMAGE_WORKFLOW):
        slide3.shapes.add_picture(IMAGE_WORKFLOW, Inches(7.5), Inches(1.8), width=Inches(5.0), height=Inches(4.5))
        
    # -------------------------------------------------------------
    # SLIDE 4: OUR PLATFORM FOUNDATIONS (SPLIT)
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    add_title(slide4, "Product Foundations: Ready Systems, Faster Delivery")
    
    bullets_s4 = [
        "AppGenius Logistics CRM & Automation:",
        "Tailored portal integrating inventory, shipping alerts, and ticket systems directly into a clean backend workflow.",
        "",
        "DataViz Analytics Dashboard:",
        "An executive-ready dashboard engine offering rapid charts, data consolidation, and drill-down metrics capabilities.",
        "",
        "VisionPro Object Detection Platform:",
        "A computer vision platform transforming camera feeds and snapshots into audit-ready operations metrics and zone compliance alerts.",
        "",
        "Enterprise Systems Manager (ESM):",
        "Central command center monitoring cloud health, logs, and database metrics in real-time."
    ]
    add_bullets(slide4, bullets_s4, width=6.2)
    
    if os.path.exists(IMAGE_DASHBOARD):
        slide4.shapes.add_picture(IMAGE_DASHBOARD, Inches(7.5), Inches(1.8), width=Inches(5.0), height=Inches(4.5))
        
    # -------------------------------------------------------------
    # SLIDE 5: LIVE SHOWCASE (PROVEN IN PRODUCTION)
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    add_title(slide5, "Proven In Production: Our Active Live Systems")
    
    bullets_s5 = [
        "We maintain fully operational digital platforms running under active workloads. You can explore these live systems on our website:",
        "",
        "Eagle Business AI Workspace: A secure, authenticated corporate workspace providing integrated AI-assisted tools and secure internal intelligence access.",
        "",
        "TN Election Insights: A responsive, high-performance feedback and constituency voter survey analytics dashboard.",
        "",
        "Aayus Naturals Devotional E-Commerce: A robust, full-featured retail platform showcasing catalog search, cart architecture, and rapid WhatsApp touchpoint sync.",
        "",
        "Lakshmi Jewellery Storefront: A highly aesthetic, fluid media storefront focused on premium branding, smooth CSS grids, and high-fidelity mobile browsing.",
        "",
        "PV School Parent Portal: A structured education portal managing academic information, admissions queries, and centralized parent communication."
    ]
    add_bullets(slide5, bullets_s5, width=11.5)
    
    # -------------------------------------------------------------
    # SLIDE 6: CASE STUDIES
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    add_title(slide6, "Client Success Stories: Real Outcomes")
    
    bullets_s6 = [
        "Case Study 1: Logistics CRM & Workflow Automation",
        "  - The Problem: A logistics team spent 14 hours/week copy-pasting tracking data.",
        "  - The Solution: Custom Python synchronization scripts that automated the entire dispatch and customer notifications pipeline.",
        "  - The Result: Reduced manual workflow effort by 95% and achieved 100% data integrity.",
        "",
        "Case Study 2: OpenSearch Energy Analytics Command Center",
        "  - The Problem: A scaling energy provider had fragmented analytics across multiple servers, making real-time KPI analysis impossible.",
        "  - The Solution: A centralized OpenSearch metrics connector generating sub-second grid health analysis and interactive trend forecasts.",
        "  - The Result: Eliminated reporting latency and reduced operational response times."
    ]
    add_bullets(slide6, bullets_s6, width=11.5)
    
    # -------------------------------------------------------------
    # SLIDE 7: ENGAGEMENT & PRICING
    # -------------------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    add_title(slide7, "Transparent Pricing & High-ROI Packages")
    
    bullets_s7 = [
        "Basic Business Website (Starting from INR 5,000)",
        "  - Responsive showcase, custom contact forms, essential SEO, and 1 year maintenance.",
        "",
        "Professional Platform Storefront (Starting from INR 8,999)",
        "  - Custom visual layouts, lead conversion paths, high-fidelity galleries, and support.",
        "",
        "Custom B2B CRM & Dashboard (Starting from INR 12,000)",
        "  - Full lead/customer database, secure follow-up tracking, dynamic reports, and 1 year support.",
        "",
        "AI & Intelligent Custom App (Starting from INR 25,000)",
        "  - AI-assisted workflow modules, advanced automation tasks, database scaling, and maintenance.",
        "",
        "*Note: Every system we build includes 1 full year of dedicated maintenance and live deployment support."
    ]
    add_bullets(slide7, bullets_s7, width=11.5)
    
    # -------------------------------------------------------------
    # SLIDE 8: THE TECHNICAL ARCHITECTURE
    # -------------------------------------------------------------
    slide8 = prs.slides.add_slide(blank_layout)
    add_title(slide8, "Robust & Scalable Technology Stack")
    
    bullets_s8 = [
        "Backend & Web Foundations: Python, Flask, Node.js, and secure RESTful API architectures.",
        "",
        "Deep Learning & Advanced Modeling: TensorFlow, Keras (LSTM forecasting models), and PyTorch ( biosignal convolutional neural networks).",
        "",
        "Databases & Search Engines: PostgreSQL, SQLite, and OpenSearch / Elasticsearch clustering.",
        "",
        "Modern Interface & Aesthetics: Responsive front-end layouts, glassmorphism UI components, fluid CSS transitions, and unified, solid-color brand blend layouts.",
        "",
        "Centralized Scalability: Single-file config system (site-config.js) dynamically updating external integrations, WhatsApp gateways, and booking endpoints across all website sub-pages instantly."
    ]
    add_bullets(slide8, bullets_s8, width=11.5)
    
    # -------------------------------------------------------------
    # SLIDE 9: START THE CONVERSATION
    # -------------------------------------------------------------
    slide9 = prs.slides.add_slide(blank_layout)
    if os.path.exists(IMAGE_COVER):
        slide9.shapes.add_picture(IMAGE_COVER, 0, 0, width=Inches(13.33), height=Inches(7.5))
        
    # Overlay Contact details on Dark Slide
    txBox = slide9.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.3), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p_badge9 = tf.paragraphs[0]
    p_badge9.text = "PARTNER WITH EAGLE INNOVATIONS"
    p_badge9.font.size = Pt(14)
    p_badge9.font.bold = True
    p_badge9.font.color.rgb = COLOR_GOLD
    p_badge9.alignment = PP_ALIGN.LEFT
    
    p_title9 = tf.add_paragraph()
    p_title9.text = "Let's build your next digital system."
    p_title9.font.size = Pt(44)
    p_title9.font.bold = True
    p_title9.font.color.rgb = COLOR_WHITE
    p_title9.alignment = PP_ALIGN.LEFT
    
    p_desc = tf.add_paragraph()
    p_desc.text = "Tell us what is slowing your team down or what you want to launch next. We will respond with a practical, high-ROI workflow design."
    p_desc.font.size = Pt(16)
    p_desc.font.color.rgb = COLOR_LIGHT_GRAY
    p_desc.alignment = PP_ALIGN.LEFT
    
    tf.add_paragraph().text = "" # spacer
    
    p_contact1 = tf.add_paragraph()
    p_contact1.text = "🟢 Direct WhatsApp / Call: +91 82204 88716"
    p_contact1.font.size = Pt(18)
    p_contact1.font.bold = True
    p_contact1.font.color.rgb = COLOR_WHITE
    
    p_contact2 = tf.add_paragraph()
    p_contact2.text = "✉️ Email Us: eaglegroupit@gmail.com"
    p_contact2.font.size = Pt(18)
    p_contact2.font.bold = True
    p_contact2.font.color.rgb = COLOR_WHITE
    
    p_contact3 = tf.add_paragraph()
    p_contact3.text = "🌐 Active Digital Workspace: https://eagleinnovations.in"
    p_contact3.font.size = Pt(18)
    p_contact3.font.bold = True
    p_contact3.font.color.rgb = COLOR_GOLD
    
    output_path = "Eagle_Innovations_Corporate_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    main()
